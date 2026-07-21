"""Programmatic graders: reopen the real files / world state and check facts.

Every grader returns (score, checks) where checks is a list of
(description, passed_bool). Score = fraction of checks passed.
No LLM is involved in grading.
"""
import os
import re

from openpyxl import load_workbook
from pptx import Presentation


def _find_file(files_dir, stem, ext):
    """Find a generated file by case-insensitive stem match."""
    if not os.path.isdir(files_dir):
        return None
    for f in os.listdir(files_dir):
        if f.lower().endswith(ext) and stem.lower() in f.lower():
            return os.path.join(files_dir, f)
    return None


def pptx_slides(path):
    """Return [(title_text, [bullet, ...]), ...] for a pptx file."""
    prs = Presentation(path)
    out = []
    for slide in prs.slides:
        title = ""
        bullets = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            texts = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
            if shape == slide.shapes.title:
                title = " ".join(texts)
            else:
                bullets.extend(texts)
        out.append((title, bullets))
    return out


def pptx_all_text(path):
    parts = []
    for title, bullets in pptx_slides(path):
        parts.append(title)
        parts.extend(bullets)
    return " \n".join(parts)


def _norm_num_text(text):
    return re.sub(r"[,\s$]", "", text.lower())


def xlsx_cells(path):
    wb = load_workbook(path)
    ws = wb.worksheets[0]
    return [list(row) for row in ws.iter_rows(values_only=True)]


def _cell_number(value, rows=None):
    """Interpret a cell as a number; evaluates simple =SUM(A1:B9) formulas."""
    if isinstance(value, (int, float)):
        return float(value)
    if isinstance(value, str):
        s = value.strip().replace("$", "").replace(",", "")
        m = re.match(r"^=SUM\(([A-Z]+)(\d+):([A-Z]+)(\d+)\)$", s, re.I)
        if m and rows is not None:
            col1 = ord(m.group(1).upper()) - 65
            col2 = ord(m.group(3).upper()) - 65
            r1, r2 = int(m.group(2)) - 1, int(m.group(4)) - 1
            total = 0.0
            for r in range(min(r1, r2), max(r1, r2) + 1):
                for c in range(min(col1, col2), max(col1, col2) + 1):
                    if r < len(rows) and c < len(rows[r]):
                        v = rows[r][c]
                        if isinstance(v, (int, float)):
                            total += float(v)
                        elif isinstance(v, str):
                            try:
                                total += float(v.replace("$", "").replace(",", ""))
                            except ValueError:
                                pass
            return total
        try:
            return float(s)
        except ValueError:
            return None
    return None


def sheet_numbers(rows):
    nums = set()
    for r in rows:
        for v in r:
            n = _cell_number(v, rows)
            if n is not None:
                nums.add(round(n, 2))
    return nums


def sheet_text(rows):
    return " ".join(str(v).lower() for r in rows for v in r if v is not None)


def find_event(world, **conds):
    """Find a calendar event added during the episode matching all conditions."""
    seeded = {"c1", "c2", "c3", "c4", "c5", "c6", "c7"}
    for e in world.events:
        if e["id"] in seeded:
            continue
        ok = True
        for k, v in conds.items():
            if k == "title_has":
                ok = ok and v.lower() in e["title"].lower()
            elif k in e:
                ok = ok and e[k] == v
            else:
                ok = False
        if ok:
            return e
    return None


def new_events(world):
    seeded = {"c1", "c2", "c3", "c4", "c5", "c6", "c7"}
    return [e for e in world.events if e["id"] not in seeded]


def minutes(t):
    h, m = t.split(":")
    return int(h) * 60 + int(m)


def overlaps(ev, others):
    s, e = minutes(ev["start"]), minutes(ev["end"])
    for o in others:
        if o["date"] != ev["date"]:
            continue
        os_, oe = minutes(o["start"]), minutes(o["end"])
        if s < oe and os_ < e:
            return True
    return False
