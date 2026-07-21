"""Real PowerPoint and Excel file creation (python-pptx / openpyxl).

The files written here are opened again by the graders (and by the human
reviewer in actual PowerPoint/Excel), so this is genuine capability, not
a simulation.
"""
import os

from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Pt

from .world import ToolError


def _resolve(files_dir, filename, ext):
    if not filename or not isinstance(filename, str):
        raise ToolError(f"'filename' must be a string ending in {ext}")
    name = os.path.basename(filename.strip())
    if not name.lower().endswith(ext):
        name += ext
    return os.path.join(files_dir, name), name


def create_presentation(files_dir, filename, slides):
    path, name = _resolve(files_dir, filename, ".pptx")
    if not isinstance(slides, list) or not slides:
        raise ToolError("'slides' must be a non-empty list of objects like "
                        '{"title": "...", "bullets": ["...", "..."]}')
    prs = Presentation()
    for i, s in enumerate(slides):
        if not isinstance(s, dict) or "title" not in s:
            raise ToolError(f"slide {i + 1} must be an object with a 'title' key "
                            f"(and optional 'bullets' list), got {s!r}")
        bullets = s.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        if not isinstance(bullets, list):
            raise ToolError(f"slide {i + 1}: 'bullets' must be a list of strings")
        bullets = [str(b) for b in bullets]
        if i == 0 and not bullets:
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # title slide
            slide.shapes.title.text = str(s["title"])
            if s.get("subtitle") and len(slide.placeholders) > 1:
                slide.placeholders[1].text = str(s["subtitle"])
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
            slide.shapes.title.text = str(s["title"])
            body = slide.placeholders[1].text_frame
            for j, b in enumerate(bullets):
                para = body.paragraphs[0] if j == 0 else body.add_paragraph()
                para.text = b
                para.font.size = Pt(20)
    prs.save(path)
    return f"created {name} with {len(slides)} slide(s)"


def create_spreadsheet(files_dir, filename, rows, sheet_name=None):
    path, name = _resolve(files_dir, filename, ".xlsx")
    if not isinstance(rows, list) or not rows or not all(isinstance(r, list) for r in rows):
        raise ToolError("'rows' must be a non-empty list of row lists, e.g. "
                        '[["Item", "Cost"], ["Laptops", 3200]]')
    wb = Workbook()
    ws = wb.active
    if sheet_name:
        ws.title = str(sheet_name)
    for r in rows:
        ws.append(r)  # strings starting with '=' become real formulas
    wb.save(path)
    return f"created {name} with {len(rows)} row(s)"


def read_spreadsheet(files_dir, filename):
    path, name = _resolve(files_dir, filename, ".xlsx")
    if not os.path.exists(path):
        raise ToolError(f"no spreadsheet named {name} exists yet")
    wb = load_workbook(path)
    out = []
    for ws in wb.worksheets:
        rows = [[c for c in row] for row in ws.iter_rows(values_only=True)]
        out.append({"sheet": ws.title, "rows": rows})
    return out
