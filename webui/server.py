"""Agent Lab — a local web console for the per-model agents.

Pick a model, type a task, press Run, and watch the loop work: the plan, each
model call streaming token by token, every tool call with the arguments the
harness actually sent, and the agent's folder (inbox, calendar, files, memory)
updating as it changes.

    python -m webui.server            then open http://127.0.0.1:8765

Binds loopback only. One run at a time, in a subprocess (webui/runner.py), so
Stop always works and the process-global harness switches can't collide.
"""
import http.server
import json
import mimetypes
import os
import queue
import shutil
import socket
import subprocess
import sys
import threading
import time
import urllib.parse
import webbrowser

import requests

HERE = os.path.dirname(os.path.abspath(__file__))
PROJECT = os.path.dirname(HERE)
STATIC = os.path.join(HERE, "static")
AGENTS_DIR = os.path.join(PROJECT, "agents")
OLLAMA_URL = "http://127.0.0.1:11434"
DEFAULT_PORT = 8765

# Rough per-size guidance for the picker; the machine, not the harness, decides.
SPEED_HINT = {
    "1b": ("instant", "Fast enough to feel live. Makes the most mistakes — the best "
                      "place to watch the harness repair a call."),
    "3b": ("quick", "A few seconds per step. The sweet spot for demos."),
    "8b": ("steady", "Tens of seconds per step. Noticeably more reliable."),
    "14b": ("slow", "Heavy on CPU-only machines. Strong tool discipline."),
    "32b": ("very slow", "Minutes per step. Fits in 32 GB RAM, tests your patience."),
}

PRESET_TASKS = [
    "Summarize my Wednesday meetings and message Jordan with the list",
    "Find a free hour on Thursday and book it as Deep work",
    "Turn Dana's Q3 sales numbers into a PowerPoint deck",
    "Build a spreadsheet of my July receipts with a total",
    "Reply to Mia about the Northwind kickoff and add it to my calendar",
    "Remember that I prefer meetings after 14:00 and never on Fridays",
]


# ----------------------------------------------------------------- agents ----

def agent_folders():
    if not os.path.isdir(AGENTS_DIR):
        return []
    out = []
    for name in sorted(os.listdir(AGENTS_DIR), key=lambda n: (len(n), n)):
        if name.startswith("_") or name.startswith("."):
            continue
        cfg_path = os.path.join(AGENTS_DIR, name, "config.json")
        if os.path.isfile(cfg_path):
            out.append(name)
    return out


def read_config(agent):
    with open(os.path.join(AGENTS_DIR, agent, "config.json"), encoding="utf-8-sig") as f:
        return json.load(f)


def agent_dir(agent):
    """Resolve an agent id to its folder, refusing anything else."""
    if agent not in agent_folders():
        raise ValueError(f"unknown agent {agent!r}")
    return os.path.join(AGENTS_DIR, agent)


def installed_tags():
    try:
        r = requests.get(f"{OLLAMA_URL}/api/tags", timeout=3)
        r.raise_for_status()
        return {m["name"]: m.get("size", 0) for m in r.json().get("models", [])}
    except Exception:
        return None  # None = server unreachable, {} = up with no models


def tag_installed(tag, tags):
    """Ollama treats llama3.1:8b and llama3.1:latest as different tags even when
    they share blobs, so only an exact match counts — except that a bare name
    means :latest."""
    if not tags:
        return False
    return tag in tags or (":" not in tag and f"{tag}:latest" in tags)


def agent_list():
    tags = installed_tags()
    out = []
    for name in agent_folders():
        cfg = read_config(name)
        folder = os.path.join(AGENTS_DIR, name)
        files_dir = os.path.join(folder, "workspace", "files")
        logs_dir = os.path.join(folder, "logs")
        mem_path = os.path.join(folder, "memory", "memory.jsonl")
        speed, blurb = SPEED_HINT.get(name, ("", ""))
        out.append({
            "id": name,
            "name": cfg.get("name", name),
            "model": cfg["model"],
            "note": cfg.get("note", ""),
            "speed": speed,
            "blurb": blurb,
            "installed": tag_installed(cfg["model"], tags),
            "files": len(os.listdir(files_dir)) if os.path.isdir(files_dir) else 0,
            "runs": len([f for f in os.listdir(logs_dir) if f.startswith("run_")])
                    if os.path.isdir(logs_dir) else 0,
            "memories": sum(1 for _ in open(mem_path, encoding="utf-8"))
                        if os.path.isfile(mem_path) else 0,
        })
    return {"agents": out, "ollama": tags is not None, "presets": PRESET_TASKS,
            "project": PROJECT}


# -------------------------------------------------------------- workspace ----

def workspace(agent):
    """The agent's folder as the browser shows it — same shape the runner emits
    during a run, so the panel renders identically live and at rest."""
    folder = agent_dir(agent)
    state_path = os.path.join(folder, "workspace", "state.json")
    files_dir = os.path.join(folder, "workspace", "files")
    mem_path = os.path.join(folder, "memory", "memory.jsonl")
    logs_dir = os.path.join(folder, "logs")

    state = {}
    if os.path.isfile(state_path):
        try:
            with open(state_path, encoding="utf-8") as f:
                state = json.load(f)
        except ValueError:
            state = {}
    if not state:  # never run: show the fixtures it will start from
        from harness.world import CALENDAR, EMAILS
        state = {"emails": [dict(e) for e in EMAILS], "events": [dict(e) for e in CALENDAR]}

    files = []
    if os.path.isdir(files_dir):
        for name in sorted(os.listdir(files_dir)):
            path = os.path.join(files_dir, name)
            if os.path.isfile(path):
                st = os.stat(path)
                files.append({"name": name, "size": st.st_size, "mtime": st.st_mtime})

    memory = []
    if os.path.isfile(mem_path):
        with open(mem_path, encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if line:
                    try:
                        memory.append(json.loads(line)["fact"])
                    except (ValueError, KeyError):
                        pass

    logs = []
    if os.path.isdir(logs_dir):
        for name in sorted(os.listdir(logs_dir), reverse=True):
            if name.startswith("run_") and name.endswith(".json"):
                logs.append({"name": name,
                             "mtime": os.path.getmtime(os.path.join(logs_dir, name))})

    return {
        "emails": state.get("emails", []),
        "sent": state.get("sent_emails", []),
        "events": sorted(state.get("events", []), key=lambda e: (e["date"], e["start"])),
        "messages": state.get("messages", []),
        "reminders": state.get("reminders", []),
        "files": files,
        "memory": memory,
        "logs": logs[:25],
        "folder": folder,
    }


def workspace_file(agent, name):
    files_dir = os.path.join(agent_dir(agent), "workspace", "files")
    path = os.path.abspath(os.path.join(files_dir, os.path.basename(str(name))))
    if os.path.dirname(path) != os.path.abspath(files_dir) or not os.path.isfile(path):
        raise ValueError(f"no such file {name!r}")
    return path


def preview(agent, name):
    """Render a generated file in the browser instead of making the user open
    PowerPoint — the whole point is to see what the agent produced."""
    path = workspace_file(agent, name)
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pptx":
        from pptx import Presentation
        slides = []
        for slide in Presentation(path).slides:
            title, body = "", []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                lines = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
                if shape == slide.shapes.title:
                    title = " ".join(lines)
                else:
                    body += lines
            slides.append({"title": title, "bullets": body})
        return {"kind": "pptx", "name": os.path.basename(path), "slides": slides}
    if ext == ".xlsx":
        from openpyxl import load_workbook
        sheets = []
        for ws in load_workbook(path, data_only=False).worksheets:
            rows = [["" if c is None else c for c in row]
                    for row in ws.iter_rows(values_only=True)]
            sheets.append({"sheet": ws.title, "rows": rows[:200]})
        return {"kind": "xlsx", "name": os.path.basename(path), "sheets": sheets}
    with open(path, "rb") as f:
        blob = f.read(20000)
    if b"\x00" in blob[:2000]:
        return {"kind": "binary", "name": os.path.basename(path),
                "size": os.path.getsize(path)}
    return {"kind": "text", "name": os.path.basename(path),
            "text": blob.decode("utf-8", errors="replace")}


def reset_agent(agent, what):
    folder = agent_dir(agent)
    done = []
    targets = {
        "world": os.path.join(folder, "workspace", "state.json"),
        "memory": os.path.join(folder, "memory", "memory.jsonl"),
    }
    for key, path in targets.items():
        if key in what and os.path.isfile(path):
            os.remove(path)
            done.append(key)
    for key, path in (("files", os.path.join(folder, "workspace", "files")),
                      ("logs", os.path.join(folder, "logs"))):
        if key in what and os.path.isdir(path):
            shutil.rmtree(path)
            os.makedirs(path, exist_ok=True)
            done.append(key)
    return done


def reveal(path):
    if sys.platform == "darwin":
        subprocess.Popen(["open", path])
    elif sys.platform.startswith("win"):
        os.startfile(path)  # noqa: S606 - local dev convenience
    else:
        subprocess.Popen(["xdg-open", path])


# -------------------------------------------------------------------- runs ----

class Run:
    """One agent subprocess, its event log, and everyone watching it."""

    def __init__(self, rid, agent, task, proc, options):
        self.id = rid
        self.agent = agent
        self.task = task
        self.proc = proc
        self.options = options
        self.started = time.time()
        self.events = []
        self.subs = []
        self.status = "running"
        self.lock = threading.Lock()

    def add(self, event):
        with self.lock:
            self.events.append(event)
            item = (len(self.events) - 1, event)
            subs = list(self.subs)
        for q in subs:
            q.put(item)

    def subscribe(self, after=-1):
        """Register a watcher and hand back everything it has not seen. `after`
        comes from Last-Event-ID, so a reconnect resumes instead of replaying."""
        q = queue.Queue()
        with self.lock:
            backlog = list(enumerate(self.events))[after + 1:]
            self.subs.append(q)
        return q, backlog

    def unsubscribe(self, q):
        with self.lock:
            if q in self.subs:
                self.subs.remove(q)

    def answer(self, cid, allow):
        try:
            self.proc.stdin.write(json.dumps({"id": cid, "allow": bool(allow)}) + "\n")
            self.proc.stdin.flush()
            return True
        except (OSError, ValueError):
            return False

    def stop(self):
        if self.proc.poll() is None:
            self.proc.terminate()
            self.status = "stopped"


class Runs:
    def __init__(self):
        self.current = None
        self.next_id = 1
        self.lock = threading.Lock()

    def start(self, agent, task, options):
        with self.lock:
            if self.current and self.current.proc.poll() is None:
                raise RuntimeError(f"{self.current.agent} is already running — "
                                   "stop it first (one model at a time).")
            cmd = [sys.executable, "-u", "-m", "webui.runner",
                   "--agent", agent, "--task", task]
            if options.get("root"):
                cmd += ["--root", options["root"]]
            if options.get("shell"):
                cmd.append("--shell")
            if options.get("yolo"):
                cmd.append("--yolo")
            if options.get("with_office"):
                cmd.append("--with-office")
            if options.get("tiers"):
                cmd.append("--tiers")
            if options.get("small"):
                cmd += ["--small", options["small"]]
            if options.get("deep"):
                cmd += ["--deep", options["deep"]]
            if options.get("max_calls"):
                cmd += ["--max-calls", str(int(options["max_calls"]))]
            env = dict(os.environ, PYTHONUNBUFFERED="1", PYTHONIOENCODING="utf-8")
            proc = subprocess.Popen(cmd, cwd=PROJECT, env=env, text=True,
                                    encoding="utf-8", errors="replace", bufsize=1,
                                    stdin=subprocess.PIPE, stdout=subprocess.PIPE,
                                    stderr=subprocess.PIPE)
            run = Run(self.next_id, agent, task, proc, options)
            self.next_id += 1
            self.current = run
        threading.Thread(target=self._pump, args=(run,), daemon=True).start()
        return run

    def _pump(self, run):
        stderr = []
        threading.Thread(target=lambda: stderr.extend(run.proc.stderr),
                         daemon=True).start()
        for line in run.proc.stdout:
            line = line.strip()
            if not line:
                continue
            try:
                run.add(json.loads(line))
            except ValueError:
                run.add({"t": "stdout", "text": line})
        code = run.proc.wait()
        time.sleep(0.05)  # let the stderr reader drain
        if run.status == "running":
            run.status = "finished" if code == 0 else "failed"
        if code not in (0, -15) and run.status != "stopped":
            tail = "".join(stderr)[-1500:].strip()
            run.add({"t": "error", "message": f"the run exited with code {code}",
                     "trace": tail})
        run.add({"t": "closed", "status": run.status, "code": code})


RUNS = Runs()


# ------------------------------------------------------------------ server ----

class Handler(http.server.BaseHTTPRequestHandler):
    protocol_version = "HTTP/1.1"
    server_version = "AgentLab"

    def log_message(self, *args):
        pass  # the console belongs to the run banner, not to request noise

    # ---- helpers ----
    def send_json(self, obj, status=200):
        body = json.dumps(obj, ensure_ascii=False, default=str).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.send_header("Cache-Control", "no-store")
        self.end_headers()
        self.wfile.write(body)

    def send_bytes(self, blob, ctype, extra=()):
        self.send_response(200)
        self.send_header("Content-Type", ctype)
        self.send_header("Content-Length", str(len(blob)))
        self.send_header("Cache-Control", "no-store")
        for k, v in extra:
            self.send_header(k, v)
        self.end_headers()
        self.wfile.write(blob)

    def body_json(self):
        length = int(self.headers.get("Content-Length") or 0)
        if not length:
            return {}
        return json.loads(self.rfile.read(length).decode("utf-8"))

    def query(self):
        parts = urllib.parse.urlparse(self.path)
        return parts.path, {k: v[0] for k, v in urllib.parse.parse_qs(parts.query).items()}

    # ---- GET ----
    def do_GET(self):
        path, q = self.query()
        try:
            if path in ("/", "/index.html"):
                return self.static_file("index.html")
            if path.startswith("/static/"):
                return self.static_file(path[len("/static/"):])
            if path == "/api/agents":
                return self.send_json(agent_list())
            if path == "/api/workspace":
                return self.send_json(workspace(q.get("agent", "")))
            if path == "/api/preview":
                return self.send_json(preview(q.get("agent", ""), q.get("name", "")))
            if path == "/api/download":
                fpath = workspace_file(q.get("agent", ""), q.get("name", ""))
                ctype = mimetypes.guess_type(fpath)[0] or "application/octet-stream"
                with open(fpath, "rb") as f:
                    blob = f.read()
                return self.send_bytes(blob, ctype, extra=[
                    ("Content-Disposition",
                     f'attachment; filename="{os.path.basename(fpath)}"')])
            if path == "/api/log":
                folder = agent_dir(q.get("agent", ""))
                name = os.path.basename(q.get("name", ""))
                with open(os.path.join(folder, "logs", name), encoding="utf-8") as f:
                    return self.send_json(json.load(f))
            if path == "/api/status":
                run = RUNS.current
                return self.send_json({"run": run.id if run else None,
                                       "agent": run.agent if run else None,
                                       "status": run.status if run else "idle"})
            if path == "/api/events":
                return self.stream_events(q)
            if path == "/api/pull":
                return self.stream_pull(q.get("model", ""))
        except ValueError as e:
            return self.send_json({"error": str(e)}, 400)
        except FileNotFoundError as e:
            return self.send_json({"error": str(e)}, 404)
        except Exception as e:  # a broken panel shouldn't take the server down
            return self.send_json({"error": f"{type(e).__name__}: {e}"}, 500)
        self.send_json({"error": "not found"}, 404)

    # ---- POST ----
    def do_POST(self):
        path, _ = self.query()
        try:
            body = self.body_json()
            if path == "/api/run":
                agent = body.get("agent", "")
                task = (body.get("task") or "").strip()
                agent_dir(agent)  # validates
                if not task:
                    raise ValueError("give the agent a task first")
                root = (body.get("root") or "").strip()
                if root and not os.path.isdir(os.path.expanduser(root)):
                    raise ValueError(f"working folder {root} does not exist")
                options = {"root": os.path.expanduser(root) if root else None,
                           "shell": body.get("shell"), "yolo": body.get("yolo"),
                           "with_office": body.get("with_office"),
                           "tiers": body.get("tiers"), "small": body.get("small"),
                           "deep": body.get("deep"), "max_calls": body.get("max_calls")}
                run = RUNS.start(agent, task, options)
                return self.send_json({"run": run.id, "agent": agent})
            if path == "/api/stop":
                run = RUNS.current
                if run:
                    run.stop()
                return self.send_json({"ok": True})
            if path == "/api/confirm":
                run = RUNS.current
                ok = bool(run) and run.answer(int(body.get("id", 0)), body.get("allow"))
                return self.send_json({"ok": ok})
            if path == "/api/reset":
                what = set(body.get("what") or [])
                return self.send_json({"cleared": reset_agent(body.get("agent", ""), what)})
            if path == "/api/reveal":
                agent = body.get("agent", "")
                sub = body.get("sub") or ""
                target = os.path.join(agent_dir(agent), *sub.split("/")) if sub \
                    else agent_dir(agent)
                if not os.path.exists(target):
                    raise ValueError("that folder does not exist yet")
                reveal(target)
                return self.send_json({"ok": True, "path": target})
        except RuntimeError as e:
            return self.send_json({"error": str(e)}, 409)
        except ValueError as e:
            return self.send_json({"error": str(e)}, 400)
        except Exception as e:
            return self.send_json({"error": f"{type(e).__name__}: {e}"}, 500)
        self.send_json({"error": "not found"}, 404)

    # ---- static ----
    def static_file(self, rel):
        path = os.path.abspath(os.path.join(STATIC, rel))
        if not path.startswith(os.path.abspath(STATIC)) or not os.path.isfile(path):
            return self.send_json({"error": "not found"}, 404)
        with open(path, "rb") as f:
            blob = f.read()
        ctype = mimetypes.guess_type(path)[0] or "text/plain"
        self.send_bytes(blob, f"{ctype}; charset=utf-8" if "text" in ctype
                        or "javascript" in ctype else ctype)

    # ---- SSE ----
    def open_stream(self):
        self.send_response(200)
        self.send_header("Content-Type", "text/event-stream; charset=utf-8")
        self.send_header("Cache-Control", "no-store")
        self.send_header("Connection", "close")
        self.end_headers()

    def push(self, obj, index=None):
        prefix = f"id: {index}\n" if index is not None else ""
        self.wfile.write(f"{prefix}data: {json.dumps(obj, ensure_ascii=False, default=str)}\n\n"
                         .encode("utf-8"))
        self.wfile.flush()

    def stream_events(self, q):
        run = RUNS.current
        want = q.get("run")
        if not run or (want and str(run.id) != str(want)):
            self.open_stream()
            self.push({"t": "closed", "status": "gone"})
            return
        try:
            after = int(self.headers.get("Last-Event-ID"))
        except (TypeError, ValueError):
            after = -1
        sub, backlog = run.subscribe(after)
        self.open_stream()
        self.close_connection = True
        try:
            for index, event in backlog:
                self.push(event, index)
            while True:
                try:
                    index, event = sub.get(timeout=10)
                except queue.Empty:
                    self.wfile.write(b": keepalive\n\n")
                    self.wfile.flush()
                    continue
                self.push(event, index)
                if event.get("t") == "closed":
                    break
        except (BrokenPipeError, ConnectionResetError):
            pass
        finally:
            run.unsubscribe(sub)

    def stream_pull(self, model):
        """Download a model from the picker, so 'run the 14B' never means
        leaving the page for a terminal."""
        self.open_stream()
        self.close_connection = True
        if not model:
            return self.push({"t": "error", "message": "no model given"})
        try:
            with requests.post(f"{OLLAMA_URL}/api/pull", json={"model": model},
                               stream=True, timeout=None) as r:
                r.raise_for_status()
                for line in r.iter_lines(decode_unicode=True):
                    if not line:
                        continue
                    msg = json.loads(line)
                    self.push({"t": "pull", "status": msg.get("status", ""),
                               "completed": msg.get("completed", 0),
                               "total": msg.get("total", 0),
                               "error": msg.get("error")})
            self.push({"t": "closed", "status": "done"})
        except (BrokenPipeError, ConnectionResetError):
            pass
        except Exception as e:
            try:
                self.push({"t": "error", "message": f"{type(e).__name__}: {e}"})
            except OSError:
                pass


class Server(http.server.ThreadingHTTPServer):
    daemon_threads = True
    allow_reuse_address = True


def free_port(start=DEFAULT_PORT):
    for port in range(start, start + 20):
        with socket.socket() as s:
            if s.connect_ex(("127.0.0.1", port)) != 0:
                return port
    return start


def main():
    sys.path.insert(0, PROJECT)
    port = free_port(int(os.environ.get("AGENT_LAB_PORT", DEFAULT_PORT)))
    url = f"http://127.0.0.1:{port}"
    tags = installed_tags()
    print(f"\n  Agent Lab  →  {url}")
    print(f"  project    {PROJECT}")
    if tags is None:
        print("  ollama     NOT RUNNING — start it, then reload the page")
    else:
        print(f"  ollama     up, {len(tags)} model(s) installed")
    print("  agents     " + ", ".join(agent_folders()))
    print("\n  Ctrl-C to stop.\n")
    if os.environ.get("AGENT_LAB_NO_BROWSER") != "1":
        threading.Timer(0.6, lambda: webbrowser.open(url)).start()
    try:
        Server(("127.0.0.1", port), Handler).serve_forever()
    except KeyboardInterrupt:
        if RUNS.current:
            RUNS.current.stop()
        print("\n  stopped.\n")


if __name__ == "__main__":
    main()
